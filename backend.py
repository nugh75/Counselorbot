from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from langchain.vectorstores import FAISS
from langchain.embeddings.huggingface import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from langchain_openai import ChatOpenAI
from typing import List
from pydantic import BaseModel
import os
import logging
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from dotenv import load_dotenv

# Configurazione logging
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

# Caricamento chiavi da file .env
load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
LLM_SERVER_URL = os.getenv("LLM_SERVER_URL", "https://api.openai.com/v1/chat/completions")

# Configurazione embeddings tramite HuggingFace
HUGGINGFACE_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"  # Modifica il modello se necessario

# Configurazione embeddings
embeddings = HuggingFaceEmbeddings(model_name=HUGGINGFACE_MODEL_NAME)
logger.info(f"Embeddings configurati con il modello {HUGGINGFACE_MODEL_NAME}")

# Configurazione database FAISS
UPLOAD_FOLDER = "./uploads"
DB_FOLDER = "./db"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DB_FOLDER, exist_ok=True)

try:
    db_path = "./db/my_database"
    faiss_index = FAISS.load_local(db_path, embeddings)
    logger.info(f"Database FAISS caricato da {db_path}")
except Exception as e:
    logger.error(f"Errore nel caricamento del database FAISS: {e}")
    faiss_index = None

# Inizializzazione FastAPI
app = FastAPI()

# Configurazione CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# *********************************************************************
# Funzione per gestire il modello LLM dinamicamente
def get_llm():
    try:
        llm = ChatOpenAI(model="gpt-4o-mini", openai_api_key=OPENAI_API_KEY)
        logger.info("LLM configurato correttamente con il modello gpt-4o-mini.")
        return llm
    except Exception as e:
        logger.error(f"Errore nella configurazione dell'LLM: {e}")
        raise HTTPException(status_code=500, detail="Errore nella configurazione dell'LLM.")

# *********************************************************************
# Funzione per salvare file caricati
def save_file(file: UploadFile, folder: str) -> str:
    try:
        file_path = os.path.join(folder, file.filename)
        with open(file_path, "wb") as f:
            f.write(file.file.read())
        logger.info(f"File salvato correttamente: {file_path}")
        return file_path
    except Exception as e:
        logger.error(f"Errore durante il salvataggio del file {file.filename}: {e}")
        raise HTTPException(status_code=500, detail=f"Errore nel salvataggio del file {file.filename}")

# *********************************************************************
# Funzione per estrarre contenuto dai file e dividerlo in chunk
def extract_content(file_path: str) -> List[Document]:
    try:
        if file_path.endswith(".pdf"):
            reader = PdfReader(file_path)
            text = [page.extract_text() for page in reader.pages if page.extract_text()]
        elif file_path.endswith(".docx"):
            doc = DocxDocument(file_path)
            text = [para.text for para in doc.paragraphs if para.text.strip()]
        elif file_path.endswith(".pptx"):
            ppt = Presentation(file_path)
            text = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
        elif file_path.endswith(".txt"):
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.readlines()
        else:
            raise ValueError("Formato file non supportato")

        # Suddividi il contenuto in chunk
        text_combined = "\n".join(text)
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1200, chunk_overlap=50)
        chunks = text_splitter.split_text(text_combined)

        logger.info(f"Contenuto estratto e suddiviso in {len(chunks)} chunk dal file {file_path}.")
        return [Document(page_content=chunk, metadata={"filename": os.path.basename(file_path)}) for chunk in chunks]

    except Exception as e:
        logger.error(f"Errore nell'estrazione del contenuto da {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Errore nell'estrazione del contenuto da {file_path}")

# *********************************************************************
# Funzione per recuperare contesto dal database FAISS
def retrieve_context(query: str) -> List[dict]:
    if not faiss_index:
        logger.warning("Database FAISS non caricato. Nessun contesto disponibile.")
        return []
    try:
        if not query.strip():
            logger.warning("La query fornita è vuota.")
            return []

        docs = faiss_index.similarity_search(query, k=5)
        if not docs:
            logger.warning("Nessun documento trovato per la query fornita.")
        else:
            logger.info(f"{len(docs)} documenti recuperati per la query: '{query}'.")

        return [
            {
                "content": doc.page_content,
                "source": {
                    "filename": doc.metadata.get("filename", "sconosciuto"),
                    "page_number": doc.metadata.get("page_number", "n/a")
                }
            }
            for doc in docs
        ]
    except Exception as e:
        logger.error(f"Errore durante il recupero del contesto: {e}")
        return []

# *********************************************************************
# Endpoint per il completamento della chat con RAG
class ChatRequest(BaseModel):
    messages: List[dict]
    temperature: float = 0.7

@app.post("/v1/chat/completions")
async def chat_completion(request: ChatRequest):
    try:
        user_message = request.messages[-1]["content"]
        logger.info(f"Messaggio ricevuto dall'utente: {user_message}")

        # Recupera contesto dal database FAISS
        context_with_sources = retrieve_context(user_message)
        if not context_with_sources:
            logger.warning("Nessun contesto rilevante trovato per la richiesta.")
            return {"llm_response": "Non ho trovato contesto rilevante per questa richiesta.", "context_chunks": []}

        # Prepara il contesto per il prompt
        context_text = "\n".join([
            f"• {item['content']} (Fonte: {item['source']['filename']}, Pagina: {item['source']['page_number']})"
            for item in context_with_sources
        ])
        logger.info("Contesto preparato per il prompt:\n" + context_text)

        # Integra il contesto nei messaggi
        augmented_messages = [
            {"role": "system", "content": f"Usa le seguenti informazioni per rispondere alla domanda:\n{context_text}"}
        ] + request.messages

        # Invia al modello LLM
        llm = get_llm()
        response = llm.predict_messages(messages=augmented_messages, temperature=request.temperature)
        logger.info(f"Risposta ricevuta dall'LLM: {response.content}")

        return {"llm_response": response.content, "context_chunks": context_with_sources}
    except Exception as e:
        logger.error(f"Errore nel completamento della chat: {e}")
        raise HTTPException(status_code=500, detail=f"Errore nel completamento della chat: {str(e)}")

# *********************************************************************
# Avvio del server
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)