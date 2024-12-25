from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from langchain.docstore.document import Document

def extract_file_content(file_path: str) -> str:
    """
    Estrae il contenuto di un file in base alla sua estensione.
    """
    try:
        if file_path.endswith(".pdf"):
            reader = PdfReader(file_path)
            return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif file_path.endswith(".docx"):
            doc = DocxDocument(file_path)
            return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        elif file_path.endswith(".pptx"):
            ppt = Presentation(file_path)
            slides_text = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            return "\n".join(slides_text)
        else:
            raise ValueError(f"Formato file non supportato: {file_path}")
    except Exception as e:
        raise RuntimeError(f"Errore durante l'elaborazione del file {file_path}: {e}")

def create_documents_from_file(file_path: str) -> list:
    """
    Crea documenti FAISS da un file processato.
    """
    content = extract_file_content(file_path)
    lines = content.split("\n")  # Dividi il contenuto in righe
    documents = [
        Document(page_content=line, metadata={"filename": file_path, "page_number": i + 1})
        for i, line in enumerate(lines) if line.strip()
    ]
    return documents