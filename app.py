from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from langchain.vectorstores import FAISS
from langchain.embeddings.huggingface import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from pydantic import BaseModel
from typing import List
import os
import logging
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from dotenv import load_dotenv
import requests

# Configurazione logging
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

# Caricamento variabili di ambiente
load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

# Directory Configurazione
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
STATIC_FOLDER = os.path.join(BASE_DIR, "static")
UPLOAD_FOLDER = os.path.join(BASE_DIR, "uploads")
DB_FOLDER = os.path.join(BASE_DIR, "db")
os.makedirs(STATIC_FOLDER, exist_ok=True)
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DB_FOLDER, exist_ok=True)

# Configurazione embeddings tramite HuggingFace
HUGGINGFACE_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"
embeddings = HuggingFaceEmbeddings(model_name=HUGGINGFACE_MODEL_NAME)
logger.info(f"Embeddings configurati con il modello {HUGGINGFACE_MODEL_NAME}")

# Configurazione Database FAISS
try:
    db_path = os.path.join(DB_FOLDER, "my_database")
    faiss_index = FAISS.load_local(db_path, embeddings)
    logger.info(f"Database FAISS caricato da {db_path}")
except Exception as e:
    logger.error(f"Errore nel caricamento del database FAISS: {e}")
    faiss_index = FAISS(embeddings)  # Crea un nuovo database se non esiste

# Inizializzazione FastAPI
app = FastAPI()

# Configurazione CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Monta i file statici
app.mount("/static", StaticFiles(directory=STATIC_FOLDER), name="static")

# Funzione per elaborare i file caricati
def extract_content(file_path: str) -> List[str]:
    try:
        if file_path.endswith(".pdf"):
            reader = PdfReader(file_path)
            return [page.extract_text() for page in reader.pages if page.extract_text()]
        elif file_path.endswith(".docx"):
            doc = DocxDocument(file_path)
            return [para.text for para in doc.paragraphs if para.text.strip()]
        elif file_path.endswith(".pptx"):
            ppt = Presentation(file_path)
            text = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return text
        else:
            raise ValueError("Formato file non supportato.")
    except Exception as e:
        logger.error(f"Errore durante l'elaborazione del file {file_path}: {e}")
        raise HTTPException(status_code=500, detail="Errore durante l'elaborazione del file.")

# Funzione per indicizzare file
def index_file(file_path: str, filename: str):
    content = extract_content(file_path)
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    documents = [Document(page_content=text, metadata={"filename": filename}) for text in splitter.split_text("\n".join(content))]
    faiss_index.add_documents(documents)
    logger.info(f"File indicizzato: {filename}")

# Endpoint per caricare e processare più file
@app.post("/upload-multiple-files")
async def upload_multiple_files(files: List[UploadFile] = File(...)):
    uploaded_files = []
    try:
        for file in files:
            file_path = os.path.join(UPLOAD_FOLDER, file.filename)
            if os.path.exists(file_path):
                logger.warning(f"File già esistente: {file.filename}. Skipping.")
                continue
            with open(file_path, "wb") as f:
                f.write(await file.read())
            logger.info(f"File caricato: {file.filename}")
            index_file(file_path, file.filename)
            uploaded_files.append(file.filename)
        
        # Salva il database FAISS
        faiss_index.save_local(os.path.join(DB_FOLDER, "my_database"))
        return {"message": "File caricati e indicizzati con successo.", "files": uploaded_files}
    except Exception as e:
        logger.error(f"Errore durante il caricamento multiplo: {e}")
        raise HTTPException(status_code=500, detail="Errore durante il caricamento multiplo dei file.")

# Funzione per caricare file da una cartella
def upload_files_in_folder(folder_path: str):
    files = []
    for file_name in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file_name)
        if os.path.isfile(file_path):
            files.append(("files", (file_name, open(file_path, "rb"))))

    if not files:
        logger.info("Nessun file trovato nella cartella.")
        return

    response = requests.post("http://localhost:8000/upload-multiple-files", files=files)
    if response.status_code == 200:
        logger.info("File caricati con successo:", response.json())
    else:
        logger.error("Errore durante il caricamento:", response.text)

# Route per servire index.html
@app.get("/", response_class=FileResponse)
async def serve_index():
    index_file = os.path.join(STATIC_FOLDER, "index.html")
    if not os.path.exists(index_file):
        logger.error("Il file index.html non è stato trovato.")
        raise HTTPException(status_code=404, detail="Il file index.html non è disponibile.")
    return FileResponse(index_file)

# Avvio del server o caricamento dei file
if __name__ == "__main__":
    import sys
    import uvicorn
    if len(sys.argv) > 1 and sys.argv[1] == "upload-folder":
        folder = UPLOAD_FOLDER
        upload_files_in_folder(folder)
    else:
        uvicorn.run(app, host="0.0.0.0", port=8000)