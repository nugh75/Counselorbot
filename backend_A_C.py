from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from langchain.vectorstores import FAISS
from langchain.embeddings.huggingface import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from typing import List
from pydantic import BaseModel
import os
import logging
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import requests
from dotenv import load_dotenv

# Configurazione logging
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

# Caricamento chiavi e configurazione iniziale
load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")

if not OPENAI_API_KEY or not ANTHROPIC_API_KEY:
    logger.error("Le chiavi API non sono configurate correttamente nel file .env")
    raise RuntimeError("Configurazione API non valida")

# Scegli il modello attivo decommentando uno dei seguenti
# ACTIVE_MODEL = "gpt-4"
# ACTIVE_API_URL = "https://api.openai.com/v1/chat/completions"

# ACTIVE_MODEL = "claude-v1"
# ACTIVE_API_URL = "https://api.anthropic.com/v1/complete"

ACTIVE_MODEL = "local-llama"
ACTIVE_API_URL = "http://localhost:1234/v1/chat/completions"

# Configurazione embeddings e database
embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
UPLOAD_FOLDER = "./uploads"
DB_FOLDER = "./db"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DB_FOLDER, exist_ok=True)

try:
    db_path = "./db/my_database"
    faiss_index = FAISS.load_local(db_path, embeddings)
    logger.info(f"Database FAISS caricato da {db_path}")
except Exception as e:
    logger.error(f"Errore nel caricamento del database FAISS: {e}")
    faiss_index = None

# Inizializzazione FastAPI
app = FastAPI()

# Configurazione CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# *********************************************************************
# Funzione per gestire i modelli LLM dinamicamente

def choose_model():
    if ACTIVE_MODEL == "gpt-4":
        return {"url": ACTIVE_API_URL, "headers": {"Authorization": f"Bearer {OPENAI_API_KEY}"}}
    elif ACTIVE_MODEL == "claude-v1":
        return {"url": ACTIVE_API_URL, "headers": {"x-api-key": ANTHROPIC_API_KEY}}
    elif ACTIVE_MODEL == "local-llama":
        return {"url": ACTIVE_API_URL, "headers": {"Content-Type": "application/json"}}
    else:
        raise ValueError("Modello non supportato.")

# *********************************************************************
# Funzione per salvare file caricati

def save_file(file: UploadFile, folder: str) -> str:
    try:
        file_path = os.path.join(folder, file.filename)
        with open(file_path, "wb") as f:
            f.write(file.file.read())
        return file_path
    except Exception as e:
        logger.error(f"Errore durante il salvataggio del file {file.filename}: {e}")
        raise HTTPException(status_code=500, detail=f"Errore nel salvataggio del file {file.filename}")

# *********************************************************************
# Funzione per estrarre contenuto dai file

def extract_content(file_path: str) -> List[str]:
    try:
        if file_path.endswith(".pdf"):
            reader = PdfReader(file_path)
            return [page.extract_text() for page in reader.pages if page.extract_text()]
        elif file_path.endswith(".docx"):
            doc = DocxDocument(file_path)
            return [para.text for para in doc.paragraphs if para.text.strip()]
        elif file_path.endswith(".pptx"):
            ppt = Presentation(file_path)
            content = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content.append(shape.text)
            return content
        elif file_path.endswith(".txt"):
            with open(file_path, "r", encoding="utf-8") as f:
                return f.readlines()
        else:
            raise ValueError("Formato file non supportato")
    except Exception as e:
        logger.error(f"Errore nell'estrazione del contenuto da {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Errore nell'estrazione del contenuto da {file_path}")

# *********************************************************************
# Funzione per recuperare contesto dal database FAISS

def retrieve_context(query: str) -> List[dict]:
    if not faiss_index:
        logger.warning("Database FAISS non caricato. Nessun contesto disponibile.")
        return []
    try:
        docs = faiss_index.similarity_search(query, k=5)
        logger.info(f"Documenti recuperati: {docs}")
        return [
            {
                "content": doc.page_content,
                "source": {
                    "filename": doc.metadata["filename"],
                    "page_number": doc.metadata["page_number"]
                }
            }
            for doc in docs
        ]
    except Exception as e:
        logger.error(f"Errore nel recupero del contesto: {str(e)}")
        return []

# *********************************************************************
# Endpoint per il completamento della chat

class ChatRequest(BaseModel):
    messages: List[dict]
    temperature: float = 0.7

@app.post("/v1/chat/completions")
async def chat_completion(request: ChatRequest):
    try:
        model_config = choose_model()
        data = {
            "messages": request.messages,
            "temperature": request.temperature
        }

        if ACTIVE_MODEL == "claude-v1":
            data["prompt"] = "\n".join([msg["content"] for msg in request.messages])

        response = requests.post(model_config["url"], json=data, headers=model_config["headers"])
        response.raise_for_status()

        response_data = response.json()
        bot_response = response_data.get("choices", [{}])[0].get("message", {}).get("content", "") \
            if ACTIVE_MODEL == "gpt-4" else response_data.get("completion", "")

        return {"llm_response": bot_response}
    except Exception as e:
        logger.error(f"Errore nel completamento della chat: {e}")
        raise HTTPException(status_code=500, detail=f"Errore nel completamento della chat: {e}")

# *********************************************************************
# Avvio del server

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)